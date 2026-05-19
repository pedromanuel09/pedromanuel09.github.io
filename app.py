import os
import re
import tempfile
import streamlit as st
from pptx import Presentation
import fitz

try:
    import win32com.client
    POWERPOINT_COM = True
except:
    POWERPOINT_COM = False


st.set_page_config(page_title="Generador OP", layout="wide")

PLANTILLA = "plantilla.pptx"


st.markdown("""
<style>
.block-container{
    padding-top: 1rem;
    max-width: 1500px;
}
h1{
    font-size: 42px !important;
    font-weight: 800 !important;
}
h2,h3{
    font-size: 22px !important;
}
.stTextInput > div > div > input{
    border-radius: 12px;
}
.stNumberInput > div > div > input{
    border-radius: 12px;
}
.stButton button{
    width:100%;
    height:55px;
    border-radius:12px;
    font-size:18px;
    font-weight:700;
}
</style>
""", unsafe_allow_html=True)


def limpiar_monto(valor):
    try:
        valor = str(valor)
        valor = valor.replace("S/", "")
        valor = valor.replace(",", "")
        valor = valor.replace(" ", "")
        return float(valor)
    except:
        return 0


def extraer_datos_pdf(pdf_file):
    texto = ""

    with fitz.open(stream=pdf_file.read(), filetype="pdf") as doc:
        for page in doc:
            texto += page.get_text() + "\n"

    datos = {
        "distrito": "",
        "tipo_inmueble": "",
        "area": "",
        "ratio": "",
        "valor_propiedad": "",
        "distrito_detalle": ""
    }

    def buscar(patron):
        r = re.search(patron, texto, re.IGNORECASE)
        return r.group(1).strip() if r else ""

    datos["distrito"] = buscar(r"DISTRITO\s*\n?(.+)")
    datos["tipo_inmueble"] = buscar(r"TIPO DE INMUEBLE\s*\n?(.+)")
    datos["area"] = buscar(r"ÁREA DE LA PROPIEDAD\s*\n?(.+)")
    datos["ratio"] = buscar(r"RATIO DE PRESTAMO\s*\n?(.+)")
    datos["valor_propiedad"] = buscar(r"VALOR DE PROPIEDAD\s*\n?(.+)")

    if datos["distrito"]:
        datos["distrito_detalle"] = datos["distrito"]

    return datos


def reemplazar_texto(prs, datos):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in datos.items():
                            marcador = "{{" + key + "}}"
                            if marcador in run.text:
                                run.text = run.text.replace(marcador, str(value))


def reemplazar_imagenes(prs, imagenes):
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                texto = shape.text.strip()

                if texto in imagenes and imagenes[texto]:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    shape.text_frame.clear()

                    slide.shapes.add_picture(
                        imagenes[texto],
                        left,
                        top,
                        width=width,
                        height=height
                    )


def convertir_pdf(pptx_path, pdf_path):
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    ppt = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
    ppt.SaveAs(os.path.abspath(pdf_path), 32)
    ppt.Close()
    powerpoint.Quit()


st.title("Generador de Oportunidades de Inversión")


datos_pdf_default = {
    "distrito": "San Juan de Lurigancho",
    "tipo_inmueble": "Casa",
    "area": "160 m2",
    "ratio": "13.90%",
    "valor_propiedad": "S/ 790,876.80",
    "distrito_detalle": "San Juan de Lurigancho"
}

if "datos_pdf" not in st.session_state:
    st.session_state["datos_pdf"] = {}


col1, col2 = st.columns(2, gap="large")

with col1:
    st.subheader("Datos de la operación")

    op = st.text_input("OP", "1315")

    monto = st.text_input(
        "Monto",
        "S/ 110,000.00"
    )

    meses = st.number_input(
        "Meses",
        min_value=1,
        value=36,
        step=1
    )

    modalidad = f"Cuota - {meses} meses"

    st.text_input(
        "Modalidad",
        modalidad,
        disabled=True
    )

with col2:
    st.subheader("Perfil del hipotecario")

    nombre = st.text_input("Nombre", "Pedro")
    giro = st.text_input("Giro del negocio", "Independiente")
    capacidad_pago = st.text_input("Capacidad de pago mensual", "S/ 5,508.00")

    p1, p2 = st.columns(2)

    with p1:
        motivo = st.text_input(
            "Motivo del préstamo",
            "Construcción de vivienda"
        )

    with p2:
        ocupacion = st.text_input(
            "Ocupación",
            "Contratista"
        )


st.divider()


pdf_col, inmueble_col = st.columns([1, 2], gap="large")

with pdf_col:
    st.subheader("PDF del inmueble")

    pdf_inmueble = st.file_uploader(
        "Subir PDF para autollenar",
        type=["pdf"]
    )

    if pdf_inmueble is not None:
        if st.button("ACTUALIZAR DATOS DESDE PDF"):
            datos_extraidos = extraer_datos_pdf(pdf_inmueble)

            datos_limpios = {}

            for k, v in datos_extraidos.items():
                if v:
                    datos_limpios[k] = v

            st.session_state["datos_pdf"] = datos_limpios
            st.success("Datos actualizados correctamente")

with inmueble_col:
    st.subheader("Datos del inmueble")

    d1, d2, d3 = st.columns(3)

    with d1:
        distrito = st.text_input(
            "Distrito",
            value=st.session_state["datos_pdf"].get(
                "distrito",
                datos_pdf_default["distrito"]
            )
        )

        tipo_inmueble = st.text_input(
            "Tipo de inmueble",
            value=st.session_state["datos_pdf"].get(
                "tipo_inmueble",
                datos_pdf_default["tipo_inmueble"]
            )
        )

    with d2:
        area = st.text_input(
            "Área",
            value=st.session_state["datos_pdf"].get(
                "area",
                datos_pdf_default["area"]
            )
        )

        ratio = st.text_input(
            "Ratio préstamo",
            value=st.session_state["datos_pdf"].get(
                "ratio",
                datos_pdf_default["ratio"]
            )
        )

    with d3:
        valor_propiedad = st.text_input(
            "Valor de propiedad",
            value=st.session_state["datos_pdf"].get(
                "valor_propiedad",
                datos_pdf_default["valor_propiedad"]
            )
        )

        distrito_detalle = st.text_input(
            "Distrito detalle",
            value=st.session_state["datos_pdf"].get(
                "distrito_detalle",
                datos_pdf_default["distrito_detalle"]
            )
        )


st.divider()


st.subheader("Proximidades")

px1, px2, px3 = st.columns(3)

with px1:
    prox1 = st.text_input("Proximidad 1", "Av. El Sol")

with px2:
    prox2 = st.text_input("Proximidad 2", "Mercado San Gabriel")

with px3:
    prox3 = st.text_input("Proximidad 3", "Av. Santa Rosa")


st.divider()


st.subheader("Rentabilidad")

r1, r2, r3 = st.columns(3)

with r1:
    tasa_hipotecaria = st.number_input(
        "Tasa hipotecaria mensual (%)",
        value=3.5,
        step=0.1
    )

with r2:
    tasa_inversionista = st.number_input(
        "Tasa inversionista mensual (%)",
        value=2.0,
        step=0.1
    )

with r3:
    diferencia_creditomype = tasa_hipotecaria - tasa_inversionista

    st.text_input(
        "Creditomype diferencia",
        f"{diferencia_creditomype:.1f}%",
        disabled=True
    )


monto_num = limpiar_monto(monto)

utilidad_mensual_num = monto_num * (tasa_inversionista / 100)
utilidad_total_num = utilidad_mensual_num * meses

utilidad_mensual = f"S/ {utilidad_mensual_num:,.2f}"
utilidad = f"S/ {utilidad_total_num:,.2f}"
tasa = f"{tasa_inversionista:.1f}%"

u1, u2, u3 = st.columns(3)

with u1:
    st.text_input(
        "Utilidad mensual",
        utilidad_mensual,
        disabled=True
    )

with u2:
    st.text_input(
        "Utilidad total",
        utilidad,
        disabled=True
    )

with u3:
    st.text_input(
        "Tasa inversionista",
        tasa,
        disabled=True
    )


st.divider()


st.subheader("Imágenes")

i1, i2, i3, i4 = st.columns(4)

with i1:
    mapa = st.file_uploader(
        "Mapa",
        type=["png", "jpg", "jpeg"]
    )

with i2:
    foto1 = st.file_uploader(
        "Foto 1",
        type=["png", "jpg", "jpeg"]
    )

with i3:
    foto2 = st.file_uploader(
        "Foto 2",
        type=["png", "jpg", "jpeg"]
    )

with i4:
    foto3 = st.file_uploader(
        "Foto 3",
        type=["png", "jpg", "jpeg"]
    )


st.divider()


if st.button("GENERAR OPORTUNIDAD"):

    if not os.path.exists(PLANTILLA):
        st.error("No se encontró plantilla.pptx")
        st.stop()

    temp_dir = tempfile.mkdtemp()

    try:
        datos = {
            "OP": op,
            "MONTO": monto,
            "MESES": meses,
            "MODALIDAD": modalidad,

            "NOMBRE": nombre,
            "GIRO": giro,
            "CAPACIDAD_PAGO": capacidad_pago,
            "MOTIVO_PRESTAMO": motivo,
            "OCUPACION": ocupacion,

            "DISTRITO": distrito,
            "TIPO_INMUEBLE": tipo_inmueble,
            "AREA": area,
            "RATIO": ratio,
            "VALOR_PROPIEDAD": valor_propiedad,
            "DISTRITO_DETALLE": distrito_detalle,

            "PROXIMIDAD_1": prox1,
            "PROXIMIDAD_2": prox2,
            "PROXIMIDAD_3": prox3,

            "UTILIDAD": utilidad,
            "UTILIDAD_MENSUAL": utilidad_mensual,
            "UTILIDAD_TOTAL": utilidad,
            "TASA": tasa,

            "TASA_HIPOTECARIA": f"{tasa_hipotecaria:.1f}%",
            "TASA_INVERSIONISTA": f"{tasa_inversionista:.1f}%",
            "DIFERENCIA_CREDITOMYPE": f"{diferencia_creditomype:.1f}%",
        }

        imagenes_guardadas = {}

        archivos = {
            "{{MAPA}}": mapa,
            "{{FOTO1}}": foto1,
            "{{FOTO2}}": foto2,
            "{{FOTO3}}": foto3,
        }

        for marcador, archivo in archivos.items():
            if archivo is not None:
                extension = archivo.name.split(".")[-1]

                ruta = os.path.join(
                    temp_dir,
                    marcador.replace("{", "").replace("}", "") + "." + extension
                )

                with open(ruta, "wb") as f:
                    f.write(archivo.getbuffer())

                imagenes_guardadas[marcador] = ruta

        prs = Presentation(PLANTILLA)

        reemplazar_texto(prs, datos)
        reemplazar_imagenes(prs, imagenes_guardadas)

        nombre_base = f"OP_{op}" if op else "OP_GENERADA"

        salida_pptx = os.path.join(temp_dir, nombre_base + ".pptx")
        salida_pdf = os.path.join(temp_dir, nombre_base + ".pdf")

        prs.save(salida_pptx)

        st.success("PowerPoint generado correctamente")

        with open(salida_pptx, "rb") as f:
            st.download_button(
                "Descargar PowerPoint",
                f,
                file_name=nombre_base + ".pptx"
            )

        try:
            if POWERPOINT_COM:
                convertir_pdf(salida_pptx, salida_pdf)

                with open(salida_pdf, "rb") as f:
                    st.download_button(
                        "Descargar PDF",
                        f,
                        file_name=nombre_base + ".pdf"
                    )

                st.success("PDF generado correctamente")

        except Exception as e:
            st.warning("No se pudo convertir a PDF automáticamente.")
            st.caption(f"Detalle técnico: {e}")

    except Exception as e:
        st.error(f"Ocurrió un error: {e}")